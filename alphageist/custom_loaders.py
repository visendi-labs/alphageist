from pptx import Presentation  # type: ignore
from typing import List, Optional
from langchain.docstore.document import Document
from langchain.document_loaders.base import BaseLoader


class PPTXLoader(BaseLoader):
    """Load pptx files."""
    file_path: str

    def __init__(self, file_path: str):
        """Initialize with file path."""
        self.file_path = file_path

    def load(self) -> List[Document]:
        """Load from file path"""
        docs: list[Document] = []

        ppt = Presentation(self.file_path)
        for n, slide in enumerate(ppt.slides):
            slide_texts = []
            for shape in slide.shapes:
                if hasattr(shape, 'text'):
                    slide_texts.append(shape.text)
            if slide_texts:
                text = '\n'.join(slide_texts)
                metadata = {"source": f"{self.file_path}"}
                docs.append(Document(page_content=text, metadata=metadata))
        return docs
